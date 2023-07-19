from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_UNDERLINE
from pptx.dml.color import RGBColor
from openpyxl import Workbook
from openpyxl.utils import get_column_letter
import yaml
import argparse

class PPTXCreator:
    # クラスの初期化メソッド
    def __init__(self, data, font_name='Yu Gothic UI', scale_width=3, scale_height=2.1):
        self.data = data  # 入力データ
        self.prs = Presentation()  # Presentationオブジェクトの生成
        self.prs.slide_width = Inches(16)  # スライドの幅を16インチに設定
        self.prs.slide_height = Inches(9)  # スライドの高さを9インチに設定
        self.font_name = font_name  # 使用するフォント名
        self.scale_width = scale_width  # テキストボックスの幅スケール
        self.scale_height = scale_height  # テキストボックスの高さスケール

        # Lean Canvasの各セクションの位置情報（左上の座標と幅、高さ）
        self.positions = {
            'Problem': (0, 0, scale_width, scale_height*2),
            'Solution': (scale_width, 0, scale_width, scale_height),
            'Key Metrics': (scale_width, scale_height, scale_width, scale_height),
            'Unique Value Proposition': (scale_width*2, 0, scale_width, scale_height*2),
            'Unfair Advantage': (scale_width*3, 0, scale_width, scale_height),
            'Channels': (scale_width*3, scale_height, scale_width, scale_height),
            'Customer Segments': (scale_width*4, 0, scale_width, scale_height*2),
            'Cost Structure': (0, scale_height*2, scale_width*2.5, scale_height),
            'Revenue Streams': (scale_width*2.5, scale_height*2, scale_width*2.5, scale_height),
        }

        self.section_order = {
            'Customer Segments': 1,
            'Problem': 2,
            'Unique Value Proposition': 3,
            'Solution': 4,
            'Channels': 5,
            'Revenue Streams': 6,
            'Cost Structure': 7,
            'Key Metrics': 8,
            'Unfair Advantage': 9,
        }


    # プロジェクト名のテキストボックスを作成するメソッド
    def create_project_name_box(self, slide, project):
        project_name_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.5))
        tf = project_name_box.text_frame
        p = tf.paragraphs[0]
        p.text = project['Project Name']
        p.font.size = Pt(20)
        p.font.name = self.font_name
        p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black color
        p.alignment = PP_ALIGN.LEFT

    # 日付のテキストボックスを作成するメソッド
    def create_date_box(self, slide, project):
        date_box = slide.shapes.add_textbox(Inches(13), Inches(0.5), Inches(2), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(project['Date'])
        p.font.size = Pt(20)
        p.font.name = self.font_name
        p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black color
        p.alignment = PP_ALIGN.RIGHT

    # Lean Canvasの各セクションを作成するメソッド
    def create_lean_canvas_sections(self, slide, project, h_offset=0.5, v_offset=1.3):
        for section, (left, top, width, height) in self.positions.items():
            if section in project['Lean Canvas']:
                text = "\n".join(project['Lean Canvas'][section])
                self.add_textbox(slide, left+h_offset, top+v_offset, width, height, text, section)

    # 'Use cases'セクションを作成するメソッド
    def create_use_cases_section(self, slide, project, h_offset=0.5, v_offset=1.3):
        if 'Use cases' in project and project['Use cases']:
            use_cases_text = "\n".join(project['Use cases'])
            use_cases_box = slide.shapes.add_textbox(Inches(0+h_offset), Inches(3*self.scale_height+v_offset), Inches(5*self.scale_width), Inches(self.scale_height))
            tf = use_cases_box.text_frame
            tf.text_anchor = MSO_ANCHOR.TOP
            tf.clear()

            p = tf.paragraphs[0]
            p.text = 'Use cases'
            p.font.size = Pt(16)
            p.font.name = self.font_name
            p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black color
            
            for line in use_cases_text.split('\n'):
                p = tf.add_paragraph()
                p.text = f"\u2022{line}"
                p.level = 0
                p.font.size = Pt(10)
                p.font.name = self.font_name

    # スライドを作成するメソッド
    def create_slide(self, project):
        blank_slide_layout = self.prs.slide_layouts[6]  # blank layout
        slide = self.prs.slides.add_slide(blank_slide_layout)
        
        self.create_project_name_box(slide, project)
        self.create_date_box(slide, project)
        self.create_lean_canvas_sections(slide, project)
        self.create_use_cases_section(slide, project)

        return slide

    def add_textbox(self, slide, left, top, width, height, text, title):
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.auto_size = False
        tf.text_anchor = MSO_ANCHOR.TOP
        tf.clear()

        # Add title with larger font size, green color, and underline
        p = tf.paragraphs[0]
        order_number = self.section_order.get(title, '')
        p.text = f"{order_number}. {title}"
        p.font.size = Pt(16)
        p.font.name = self.font_name
        p.font.color.rgb = RGBColor(0x00, 0x80, 0x00)  # Green color
        p.font.underline = MSO_UNDERLINE.SINGLE_LINE

        # Add content with bullet points
        for line in text.split('\n'):
            p = tf.add_paragraph()
            p.text = f"\u2022{line}"
            p.level = 0
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black color
            p.font.name = self.font_name

        # Add border to textbox
        line = textbox.line
        line.color.rgb = RGBColor(0, 0, 0)  # Black color
        line.width = Pt(1.0)

        return slide

    def create_presentation(self):
        for project in self.data['Projects']:
            if not project.get('Project Name'):
                continue
            if project.get('Project Name') is None:
                continue
            self.create_slide(project)

    def save(self, output):
        self.prs.save(output)

class XLSXCreator:
    def __init__(self, data):
        self.data = data
        self.wb = Workbook()

    def create_worksheet(self):
        ws = self.wb.active
        ws.title = "Projects"

        # ヘッダーの作成
        headers = ['Project Name', 'Date'] + list(self.data['Projects'][0]['Lean Canvas'].keys())
        ws.append(headers)

        for project in self.data['Projects']:
            if not project.get('Project Name'):
                continue
            if project.get('Project Name') is None:
                continue
            row = [project['Project Name'], project['Date']]
            for section, content in project['Lean Canvas'].items():
                row.append('\n'.join(content))
            ws.append(row)

            # Adjust column width
            for idx, column in enumerate(ws.columns, start=1):
                max_length = 0
                column = [cell for cell in column]
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(cell.value)
                    except:
                        pass
                adjusted_width = (max_length + 2)
                ws.column_dimensions[get_column_letter(idx)].width = adjusted_width

    def save(self, output):
        self.wb.save(output)


# コマンドライン引数を解析する関数です
def parse_arguments():
    parser = argparse.ArgumentParser()
    parser.add_argument('-i', '--input', default="sample.yaml", help='Input YAML file')
    parser.add_argument('-p', '--pptx-output', default="sample.pptx", help='Output file (PPTX)')
    parser.add_argument('-x', '--xlsx-output', default="sample.xlsx", help='Output file (XLSX)')
    args = parser.parse_args()
    return args

# ファイル処理を行う関数です
def process_files(input_file, pptx_output, xlsx_output):
    with open(input_file, encoding="utf-8") as file:
        data = yaml.safe_load(file)

    pptx_creator = PPTXCreator(data)
    pptx_creator.create_presentation()
    pptx_creator.save(pptx_output)

    xlsx_creator = XLSXCreator(data)
    xlsx_creator.create_worksheet()
    xlsx_creator.save(xlsx_output)

if __name__ == "__main__":
    args = parse_arguments()
    process_files(args.input, args.pptx_output, args.xlsx_output)
